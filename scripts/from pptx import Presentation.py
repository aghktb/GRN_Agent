from pptx import Presentation

# Create presentation
prs = Presentation()

# Helper function to add slide
def add_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bp in enumerate(bullet_points):
        if i == 0:
            body.text = bp
        else:
            p = body.add_paragraph()
            p.text = bp
            p.level = 1

# Slides content
add_slide(
    "Cheng Lab – Overall Role",
    [
        "Computational, AI, and bioinformatics backbone of the EDGE–NSF HSM project",
        "Integrates multi-omics data to uncover regulatory mechanisms of heat stress memory",
        "Leads deep learning, gene regulatory network modeling, and protein structure prediction"
    ]
)

add_slide(
    "Aim 1: Cis-Regulatory Element (CRE) Identification",
    [
        "Identify HSM-associated cis-regulatory elements (CREs)",
        "Promoter regions (≤2 kb upstream of TSS) and distal enhancers",
        "Integrate RNA-seq, ATAC-seq, and histone modification data",
        "Motif discovery using MEME and TomTom"
    ]
)

add_slide(
    "Aim 1: Gene Regulatory Network (GRN) Reconstruction",
    [
        "Reconstruct gene regulatory networks controlling HSM",
        "Use GNET2 (probabilistic modeling)",
        "Use GRNFormer (graph transformer deep learning)",
        "Generate consensus GRNs via cross-validation"
    ]
)

add_slide(
    "Aim 1: Deep Learning Multi-Omics Integration",
    [
        "Develop attention-based deep learning models",
        "Predict gene expression dynamics during HSM",
        "Integrate ATAC-seq, histone marks, and CREs",
        "Use attention weights to identify key regulatory features"
    ]
)

add_slide(
    "Aim 1: Protein Structure & Function Prediction",
    [
        "Predict protein tertiary and quaternary structures of HSM candidates",
        "Use MULTICOM4 (CASP16 top-performing system)",
        "Predict protein–protein interactions and functional annotations",
        "Tools: TransFun, TransFew, DeepGraphGO"
    ]
)

add_slide(
    "Aim 2: Computational Validation of CrHSF1 & CrFGT1",
    [
        "Validate predicted targets using mutant and ChIP-seq data",
        "Screen interaction partners via quaternary structure modeling",
        "Rank partners based on confidence scores",
        "Support mechanistic understanding of transgenerational HSM"
    ]
)

add_slide(
    "Aim 3: HSM Candidate Gene Prioritization",
    [
        "Integrate results from Aims 1 and 2",
        "Rank high-confidence HSM candidate genes",
        "Assess evolutionary conservation across species",
        "Support selection for functional validation"
    ]
)

add_slide(
    "Open-Source & Community Deliverables",
    [
        "Release AI and bioinformatics tools on GitHub",
        "Provide documented pipelines for GRN and multi-omics integration",
        "Disseminate tools via conferences and publications",
        "Provide user support to research community"
    ]
)

add_slide(
    "Training, Collaboration, and Broader Impacts",
    [
        "Train students in AI-driven computational biology",
        "Support cross-lab data interpretation and publications",
        "Advance reproducible and interpretable AI for biology"
    ]
)

# Save presentation
file_path = "Cheng_Lab_Deliverables_EDGE_NSF_HSM.pptx"
prs.save(file_path)

file_path